import sys
import base64
import os
import io
import re
import html

from typing import BinaryIO, Any
from operator import attrgetter

from ._html_converter import HtmlConverter
from ._llm_caption import llm_caption
from .._base_converter import DocumentConverter, DocumentConverterResult
from .._stream_info import StreamInfo
from .._exceptions import MissingDependencyException, MISSING_DEPENDENCY_MESSAGE

# Try loading optional (but in this case, required) dependencies
# Save reporting of any exceptions for later
_dependency_exc_info = None
try:
    import pptx
except ImportError:
    # Preserve the error and stack trace for later
    _dependency_exc_info = sys.exc_info()


def split_and_merge_gif(input_gif, output_png='merged_frames.png'):
    """
    Split GIF into first, middle, and last frames, then merge horizontally to PNG
    
    Args:
        input_gif: Input GIF file path
        output_png: Output PNG file path (default: 'merged_frames.png')
    
    Returns:
        output_png: Output file path
    """
    try:
        from PIL import Image
    except ImportError:
        # If PIL is not available, just return None
        return None
    
    # Open GIF file
    gif = Image.open(input_gif)
    
    # Get total frame count
    frame_count = 0
    try:
        while True:
            gif.seek(frame_count)
            frame_count += 1
    except EOFError:
        pass
    
    # Calculate middle frame index (for even numbers, take the earlier one)
    middle_index = (frame_count // 2) - 1 if frame_count % 2 == 0 else frame_count // 2
    
    # Extract three key frames
    gif.seek(0)
    first_frame = gif.convert('RGBA')
    
    gif.seek(middle_index)
    middle_frame = gif.convert('RGBA')
    
    gif.seek(frame_count - 1)
    last_frame = gif.convert('RGBA')
    
    # Create horizontally merged image
    width, height = first_frame.size
    merged_image = Image.new('RGBA', (width * 3, height))
    
    # Paste three frames
    merged_image.paste(first_frame, (0, 0))
    merged_image.paste(middle_frame, (width, 0))
    merged_image.paste(last_frame, (width * 2, 0))
    
    # Save as PNG
    merged_image.save(output_png, 'PNG')
    gif.close()
    
    return output_png


ACCEPTED_MIME_TYPE_PREFIXES = [
    "application/vnd.openxmlformats-officedocument.presentationml",
]

ACCEPTED_FILE_EXTENSIONS = [".pptx"]


class PptxConverter(DocumentConverter):
    """
    Converts PPTX files to Markdown. Supports heading, tables and images with alt text.
    """

    def __init__(self):
        super().__init__()
        self._html_converter = HtmlConverter()

    def accepts(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,  # Options to pass to the converter
    ) -> bool:
        mimetype = (stream_info.mimetype or "").lower()
        extension = (stream_info.extension or "").lower()

        if extension in ACCEPTED_FILE_EXTENSIONS:
            return True

        for prefix in ACCEPTED_MIME_TYPE_PREFIXES:
            if mimetype.startswith(prefix):
                return True

        return False

    def convert(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,  # Options to pass to the converter
    ) -> DocumentConverterResult:
        # Check the dependencies
        if _dependency_exc_info is not None:
            raise MissingDependencyException(
                MISSING_DEPENDENCY_MESSAGE.format(
                    converter=type(self).__name__,
                    extension=".pptx",
                    feature="pptx",
                )
            ) from _dependency_exc_info[
                1
            ].with_traceback(  # type: ignore[union-attr]
                _dependency_exc_info[2]
            )

        # Perform the conversion
        presentation = pptx.Presentation(file_stream)
        md_content = ""
        slide_num = 0
        image_counter = 0  # Track images across all slides
        images_dir = kwargs.get("save_images_dir")  # Optional: directory to save images
        llm_prompt_png = kwargs.get("llm_prompt")  # Optional: LLM prompt for images
        
        # LLM prompt for GIF merged images
        llm_prompt_gif = (
            "This is a composite image showing three frames from an animation "
            "(start, middle, end) arranged horizontally. "
            "Describe what this animation shows and what motion or change occurs "
            "across these three frames in Chinese. " + (llm_prompt_png or "")
        )
        for slide in presentation.slides:
            slide_num += 1

            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title

            def get_shape_content(shape, **kwargs):
                nonlocal md_content, image_counter
                # Pictures
                if self._is_picture(shape):
                    # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069

                    llm_description = ""
                    alt_text = ""

                    # Extract image metadata (needed for both LLM and non-LLM paths)
                    image_filename = shape.image.filename
                    image_extension = None
                    if image_filename:
                        image_extension = os.path.splitext(image_filename)[1]

                    # Check if this is a GIF and convert it early (for both LLM and saving)
                    content_type = shape.image.content_type or "image/png"
                    is_gif = 'gif' in content_type.lower()
                    image_blob_for_llm = shape.image.blob
                    content_type_for_llm = content_type
                    saved_path = None
                    
                    # Save all images (GIF and normal) in one place
                    if images_dir:
                        from pathlib import Path
                        images_path = Path(images_dir)
                        images_path.mkdir(parents=True, exist_ok=True)
                        
                        # Use sequential naming
                        image_counter += 1
                        
                        if is_gif:
                            # Save original GIF
                            temp_filename = f"slide{slide_num}_image{image_counter}.gif"
                            temp_path = images_path / temp_filename
                            with open(temp_path, 'wb') as f:
                                f.write(shape.image.blob)
                            
                            # Convert to merged PNG
                            merged_filename = f"slide{slide_num}_image{image_counter}_merged.png"
                            merged_path = images_path / merged_filename
                            result = split_and_merge_gif(str(temp_path), str(merged_path))
                            
                            if result is not None:
                                # Use merged PNG for both LLM and final output
                                with open(merged_path, 'rb') as f:
                                    image_blob_for_llm = f.read()
                                content_type_for_llm = "image/png"
                                saved_path = merged_path
                                promptIn = llm_prompt_gif
                            else:
                                # Conversion failed, save original as PNG
                                saved_path = images_path / f"slide{slide_num}_image{image_counter}.png"
                                with open(saved_path, 'wb') as f:
                                    f.write(shape.image.blob)
                                image_blob_for_llm = shape.image.blob
                                promptIn = llm_prompt_png
                        else:
                            # Normal image processing
                            ext = image_extension or '.png'
                            if not ext.startswith('.'):
                                ext = '.' + ext
                            saved_filename = f"slide{slide_num}_image{image_counter}{ext}"
                            saved_path = images_path / saved_filename
                            with open(saved_path, 'wb') as f:
                                f.write(shape.image.blob)
                            promptIn = llm_prompt_png


                    # Potentially generate a description using an LLM
                    llm_client = kwargs.get("llm_client")
                    llm_model = kwargs.get("llm_model")
                    if llm_client is not None and llm_model is not None:
                        # Prepare a file_stream and stream_info for the image data
                        image_stream_info = StreamInfo(
                            mimetype=content_type_for_llm,
                            extension=image_extension,
                            filename=image_filename,
                        )

                        image_stream = io.BytesIO(image_blob_for_llm)

                        # Caption the image
                        try:
                            llm_description = llm_caption(
                                image_stream,
                                image_stream_info,
                                client=llm_client,
                                model=llm_model,
                                prompt=promptIn
                                #prompt=kwargs.get("llm_prompt"),
                            )
                            print(f"[LLM Description] Slide {slide_num}, Image {image_counter}: {llm_description}")
                        except Exception as e:
                            # Unable to generate a description
                            print(f"[LLM Description] Slide {slide_num}, Image {image_counter}: Failed - {e}")
                            pass

                    # Also grab any description embedded in the deck
                    try:
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                    except Exception:
                        # Unable to get alt text
                        pass

                    # Prepare the alt, escaping any special characters
                    alt_text = "\n".join([llm_description, alt_text]) or shape.name
                    alt_text = re.sub(r"[\r\n\[\]]", " ", alt_text)
                    alt_text = re.sub(r"\s+", " ", alt_text).strip()

                    # If keep_data_uris is True, use base64 encoding for images
                    if kwargs.get("keep_data_uris", False):
                        blob = shape.image.blob
                        content_type = shape.image.content_type or "image/png"
                        b64_string = base64.b64encode(blob).decode("utf-8")
                        md_content += f"\n![{alt_text}](data:{content_type};base64,{b64_string})\n"
                    else:
                        # Use the saved path (already saved in L177 section)
                        if saved_path is not None:
                            # Use absolute path in markdown
                            filename = f"{images_path.name}/{saved_path.name}"
                        else:
                            # A placeholder name (original behavior when no images_dir)
                            filename = re.sub(r"\W", "", shape.name) + ".jpg"
                        
                        md_content += "\n![" + alt_text + "](" + filename + ")\n"

                # Tables
                if self._is_table(shape):
                    md_content += self._convert_table_to_markdown(shape.table, **kwargs)

                # Charts
                if shape.has_chart:
                    md_content += self._convert_chart_to_markdown(shape.chart)

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\n"
                    else:
                        md_content += shape.text + "\n"

                # Group Shapes
                if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
                    sorted_shapes = sorted(
                        shape.shapes,
                        key=lambda x: (
                            float("-inf") if not x.top else x.top,
                            float("-inf") if not x.left else x.left,
                        ),
                    )
                    for subshape in sorted_shapes:
                        get_shape_content(subshape, **kwargs)

            sorted_shapes = sorted(
                slide.shapes,
                key=lambda x: (
                    float("-inf") if not x.top else x.top,
                    float("-inf") if not x.left else x.left,
                ),
            )
            for shape in sorted_shapes:
                get_shape_content(shape, **kwargs)

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(markdown=md_content.strip())

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

    def _convert_table_to_markdown(self, table, **kwargs):
        # Write the table as HTML, then convert it to Markdown
        html_table = "<html><body><table>"
        first_row = True
        for row in table.rows:
            html_table += "<tr>"
            for cell in row.cells:
                if first_row:
                    html_table += "<th>" + html.escape(cell.text) + "</th>"
                else:
                    html_table += "<td>" + html.escape(cell.text) + "</td>"
            html_table += "</tr>"
            first_row = False
        html_table += "</table></body></html>"

        return (
            self._html_converter.convert_string(html_table, **kwargs).markdown.strip()
            + "\n"
        )

    def _convert_chart_to_markdown(self, chart):
        try:
            md = "\n\n### Chart"
            if chart.has_title:
                md += f": {chart.chart_title.text_frame.text}"
            md += "\n\n"
            data = []
            category_names = [c.label for c in chart.plots[0].categories]
            series_names = [s.name for s in chart.series]
            data.append(["Category"] + series_names)

            for idx, category in enumerate(category_names):
                row = [category]
                for series in chart.series:
                    row.append(series.values[idx])
                data.append(row)

            markdown_table = []
            for row in data:
                markdown_table.append("| " + " | ".join(map(str, row)) + " |")
            header = markdown_table[0]
            separator = "|" + "|".join(["---"] * len(data[0])) + "|"
            return md + "\n".join([header, separator] + markdown_table[1:])
        except ValueError as e:
            # Handle the specific error for unsupported chart types
            if "unsupported plot type" in str(e):
                return "\n\n[unsupported chart]\n\n"
        except Exception:
            # Catch any other exceptions that might occur
            return "\n\n[unsupported chart]\n\n"
